import torch
import torch.nn as nn
import torch.optim as optim
from torch.utils.data import DataLoader, TensorDataset
import numpy as np
from sklearn.preprocessing import StandardScaler
import pandas as pd
import matplotlib.pyplot as plt
import os
import io
import sys
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 1. 数据准备（修改输入列为Time）
# file_path = "./processed_data.xlsx"
file_path = "E:/Desktop/papper_PINN_SC/PINN_SC/raw_data/test_data/800V/output_files/Ids-data_Vds=800_Vgs=9.csv"
# 根据文件扩展名使用不同的读取方式
file_extension = os.path.splitext(file_path)[1].lower()
data_frames = []

if file_extension == '.xlsx':
    excel_file = pd.ExcelFile(file_path)
    for sheet in excel_file.sheet_names:
        df = pd.read_excel(excel_file, sheet_name=sheet)
        df.columns = df.columns.str.strip().str.lower()
        df = df.rename(columns={
        'time(s)': 'Time',
        'vds': 'Vds',
        'vgs': 'Vgs',
        'ids': 'Ids'
    })[['Ids', 'Time', 'Vds', 'Vgs']]  # 调整列顺序
    data_frames.append(df)
elif file_extension == '.csv':
    df = pd.read_csv(file_path)
    df.columns = df.columns.str.strip().str.lower()
    df = df.rename(columns={
    'time(s)': 'Time',
    'vds': 'Vds',
    'vgs': 'Vgs',
    'ids': 'Ids'
    })[['Ids', 'Time', 'Vds', 'Vgs']]  # 调整列顺序
    data_frames.append(df)
else:
    raise ValueError(f"不支持的文件格式: {file_extension}，仅支持.xlsx和.csv文件")

combined_df = pd.concat(data_frames, ignore_index=True)
# 按时间排序数据
#  combined_df.sort_values(by='Time', inplace=True)
raw_data = combined_df.to_numpy()
# 输入列为Time, Vds, Vgs
inputs_data = raw_data[:, 1:].copy()  # [Time, Vds, Vgs]
targets_data = raw_data[:, 0].reshape(-1, 1)

# 数据标准化
input_scaler = StandardScaler()
target_scaler = StandardScaler()
X_scaled = input_scaler.fit_transform(inputs_data)
y_scaled = target_scaler.fit_transform(targets_data)

# 保存标准化参数
input_mean = torch.tensor(input_scaler.mean_, dtype=torch.float32)
input_std = torch.tensor(input_scaler.scale_, dtype=torch.float32)
target_mean = torch.tensor(target_scaler.mean_[0], dtype=torch.float32)
target_std = torch.tensor(target_scaler.scale_[0], dtype=torch.float32)

# 创建数据集（关闭shuffle保证时间序列）
dataset = TensorDataset(torch.tensor(X_scaled, dtype=torch.float32),
                        torch.tensor(y_scaled, dtype=torch.float32))
# dataloader = DataLoader(dataset, batch_size=len(dataset), shuffle=False)  # 全量数据作为单个batch
dataloader = DataLoader(dataset, batch_size=len(dataset), shuffle=False)
# 2. 定义改进的神经网络模型（新增Epi_thickness和T_initial参数）
class PINN(nn.Module):
    def __init__(self, dropout_rate=0.02):
        super().__init__()
        self.net = nn.Sequential(
            nn.Linear(3, 128),
            nn.Tanh(),
            nn.Dropout(dropout_rate),
            nn.Linear(128, 128),
            nn.Tanh(),
            nn.Dropout(dropout_rate),
            nn.Linear(128, 128),
            nn.Tanh(),
            nn.Linear(128, 1)
        )
        # 可学习参数
        self.param1 = nn.Parameter(torch.tensor(75.132, dtype=torch.float32))
        self.param2 = nn.Parameter(torch.tensor(2.062, dtype=torch.float32))
        self.param3 = nn.Parameter(torch.tensor(6.11, dtype=torch.float32))
        self.param4 = nn.Parameter(torch.tensor(3.5389, dtype=torch.float32))
        self.param5 = nn.Parameter(torch.tensor(-0.102, dtype=torch.float32))
        self.param6 = nn.Parameter(torch.tensor(10.138, dtype=torch.float32))
        self.param7 = nn.Parameter(torch.tensor(1.222, dtype=torch.float32))
        self.param8 = nn.Parameter(torch.tensor(5.0, dtype=torch.float32))
        self.param9 = nn.Parameter(torch.tensor(-0.5, dtype=torch.float32))
        # 新增参数
        self.epi_thickness = nn.Parameter(torch.tensor(1.0e-5, dtype=torch.float32))
        self.T_initial = nn.Parameter(torch.tensor(300.0, dtype=torch.float32))

    def forward(self, x):
        return self.net(x)

model = PINN()

# 3. 物理模型定义（保持原结构，T由外部传入）
def physics_model(Vds, Vgs, T, param1, param2, param3, param4, param5, param6, param7, param8, param9):
    # 转换为张量（如果输入是标量或numpy数组）
    if not isinstance(Vds, torch.Tensor):
        Vds = torch.tensor(Vds, dtype=torch.float32)
    if not isinstance(Vgs, torch.Tensor):
        Vgs = torch.tensor(Vgs, dtype=torch.float32)
    if not isinstance(T, torch.Tensor):
        T = torch.tensor(T, dtype=torch.float32)

    # 物理公式计算
    # NET5_value = param3 * 1 / (param1*(T/300)**-param7 + param6*(T/300)**param2)  #迁移率项，受Vgs与Temperature的共同影响  stable
    NET5_value = param3 * 1 / (param1*(T/300)**-1.222 + param6*(T/300)**2.062)
    # NET5_value = param3 * 1 / (75.13*(T/300)**-1.222 + 10*(T/300)**2.062)
    # NET3_value = -0.004263*T + 3.422579  #stable
    # NET3_value = -0.00484*T + param4
    NET3_value = -0.004263*T + 3.53892
    # NET3_value = -param4*T + 3.422579
    p9 = param5
    # NET2_value = -0.005 * Vgs + 0.165
    NET2_value = -0.005 * Vgs + 0.165
    NET1_value = -0.1717 * Vgs + 3.5755
    term3 = (torch.log(1 + torch.exp(Vgs - NET3_value)))**2 - (torch.log(1 + torch.exp(Vgs - NET3_value - (NET2_value * Vds * ((1 + torch.exp(p9 * Vds))**NET1_value)))))**2
    term1 = NET5_value * (Vgs - NET3_value)
    term2 = 1 + 0.0005 * Vds
    return term2 * term1 * term3

# 4. 训练配置
optimizer = optim.AdamW([
    {'params': model.net.parameters(), 'lr': 0.001},
    {'params': [model.param1, model.param2, model.param3, model.param4, model.param5,
                model.param6, model.param7, model.param8, model.param9,
                model.epi_thickness, model.T_initial], 'lr': 0.1}
], weight_decay=1e-4)
scheduler = optim.lr_scheduler.ReduceLROnPlateau(optimizer, 'min', patience=20, factor=0.5)
mse_loss = nn.MSELoss()

# 5. 训练循环（全量数据处理）
num_epochs = 500
alpha = 0.9
w_data, w_physics = 0.5, 0.5
losses = []
data_metrics = []
physics_metrics = []

def integrate_temperature(time_raw, vds_raw, vgs_raw, model, physics_model):
    """温度积分计算函数（支持梯度传播）"""
    # 初始化温度序列（保持标量形式）
    T_list = [model.T_initial.clone().requires_grad_(True)]  # 初始温度作为可求导标量
    physics_preds = []

    delta_time = torch.zeros_like(time_raw)
    delta_time[1:] = time_raw[1:] - time_raw[:-1]

    for i in range(len(time_raw)):
        if i > 0:
            with torch.enable_grad():
                # 使用前一步温度值（保持标量形式）
                ids_prev = physics_model(vds_raw[i], vgs_raw[i], T_list[i-1].detach(),
                                      model.param1, model.param2, model.param3,
                                      model.param4, model.param5, model.param6,
                                      model.param7, model.param8, model.param9)

                # 温度计算（确保所有操作保持标量形式）
                delta_T = 2 * vds_raw[i] / 1e-5 * (ids_prev / 20e-6) * delta_time[i] / (
                           300*(5.13 - 1001 / T_list[i-1] + (3.23e4)/(T_list[i-1]**2))*3200)
                current_T = T_list[i-1] + delta_T
        else:
            # 首次迭代使用克隆的初始温度
            current_T = T_list[0].clone()

        # 将当前温度添加至列表（保持标量形式）
        T_list.append(current_T)
        # 计算物理预测值
        physics_preds.append(
            physics_model(vds_raw[i], vgs_raw[i], T_list[i],
                        model.param1, model.param2, model.param3,
                        model.param4, model.param5, model.param6,
                        model.param7, model.param8, model.param9)
        )

    # 转换时排除初始温度（索引0）
    T = torch.stack(T_list[1:])  # 所有温度均为标量形式
    physics_preds = torch.stack(physics_preds)
    return T, physics_preds

for epoch in range(num_epochs):
    model.train()
    for X_batch, y_batch in dataloader:  # 单次迭代（全量数据）
        # 数据前向传播
        pred = model(X_batch)
        loss_data = mse_loss(pred, y_batch)

        # 在训练循环中使用：
        # 计算物理损失
        time_raw = X_batch[:, 0] * input_std[0] + input_mean[0]
        vds_raw = X_batch[:, 1] * input_std[1] + input_mean[1]
        vgs_raw = X_batch[:, 2] * input_std[2] + input_mean[2]

        # 执行温度积分
        T, physics_preds = integrate_temperature(time_raw, vds_raw, vgs_raw, model, physics_model)

        # 标准化处理
        physics_pred_scaled = (physics_preds - target_mean) / target_std

        # 计算损失（增加数值稳定性）
        loss_physics = mse_loss(pred.squeeze(), physics_pred_scaled)

        # 动态权重调整
        current_data = loss_data.item()
        current_physics = loss_physics.item()
        total = current_data + current_physics + 1e-8
        new_wd = current_data / total
        new_wp = current_physics / total
        w_data = alpha * w_data + (1 - alpha) * new_wd
        w_physics = alpha * w_physics + (1 - alpha) * new_wp

        # 总损失
        total_loss = w_data * loss_data + w_physics * loss_physics

        # 反向传播
        optimizer.zero_grad()
        total_loss.backward()
        torch.nn.utils.clip_grad_norm_(model.parameters(), 1)
        optimizer.step()

    # 更新学习率和记录指标
    scheduler.step(total_loss.item())
    losses.append(total_loss.item())
    data_metrics.append(loss_data.item())
    physics_metrics.append(loss_physics.item())

    # 打印训练信息
    if (epoch + 1) % 50 == 0:
        print(f"Epoch {epoch + 1:04d} | Loss: {total_loss.item():.3e} | "
              f"Data: {loss_data.item():.2e} | Physics: {loss_physics.item():.2e} | "
              f"LR: {optimizer.param_groups[0]['lr']:.2e} | "
              f"params: {model.param1.item():.3f}, {model.param2.item():.3f}, "
              f"{model.param3.item():.3f}, {model.param4.item():.3f}, {model.param5.item():.3f}, {model.param6.item():.3f}, {model.param7.item():.3f}")

# 6. 模型预测与可视化
def neural(time, Vds, Vgs):
    input_arr = np.array([[time, Vds, Vgs]])
    scaled_input = input_scaler.transform(input_arr)
    with torch.no_grad():
        scaled_output = model(torch.tensor(scaled_input, dtype=torch.float32))
    return target_scaler.inverse_transform(scaled_output.numpy())[0][0]# def plot_results(true_values, predicted_values, physic_values):

# 7. 温度变化可视化
def plot_temperature(time, temperature):
    plt.figure(figsize=(10, 6))
    plt.plot(time, temperature, label="Temperature", color='r', linewidth=2)
    plt.xlabel("Time (s)")
    plt.ylabel("Temperature")
    plt.title("Device Temperature Evolution")
    plt.legend()
    plt.grid(True)
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight')
    plt.close()
    buf.seek(0)
    return buf

# 获取完整温度序列（使用全量数据）
with torch.no_grad():
    # 准备完整数据
    full_inputs = torch.tensor(X_scaled, dtype=torch.float32)
    full_time = full_inputs[:, 0] * input_std[0] + input_mean[0]
    full_vds = full_inputs[:, 1] * input_std[1] + input_mean[1]
    full_vgs = full_inputs[:, 2] * input_std[2] + input_mean[2]

    # 计算温度序列
    T_sequence, _ = integrate_temperature(full_time, full_vds, full_vgs, model, physics_model)

    # 获取温度图图像
    temp_buf = plot_temperature(full_time.numpy(), T_sequence.numpy())

# 对比结果可视化（使用前300个样本）
def plot_comparison(time, true, pred, physics):
    plt.figure(figsize=(10, 6))
    plt.plot(time, true, label="True Values", marker='o', markersize=3)
    plt.plot(time, pred, label="Predicted Values", marker='x', markersize=3)
    plt.plot(time, physics, label="Physics Values", marker='.', markersize=3)
    plt.xlabel("Time (s)")
    plt.ylabel("Drain Current (A)")
    plt.title("Current Comparison Over Time")
    plt.legend()
    plt.grid(True)
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight')
    plt.close()
    buf.seek(0)
    return buf
# 主可视化流程
# 修复AttributeError的核心修改部分
with torch.no_grad():
    # 2. 准备对比数据（前300个样本）
    # 正确获取时间序列（全numpy运算）
    input_std_np = input_std.numpy()
    input_mean_np = input_mean.numpy()
    partial_time = X_scaled[:, 0] * input_std_np[0] + input_mean_np[0]

    # 修复神经网络预测值的转换
    predicted_values = [neural(Vds, Vgs, T) for Vds, Vgs, T in inputs_data[:]]  # 原始输出类型

    # 正确转换物理模型输出
    # 修复部分：批量计算物理模型的预测值 确保physic_values为numpy类型
    physic_values = physics_model(full_vds, full_vgs, T_sequence,
                             model.param1, model.param2, model.param3,
                             model.param4, model.param5, model.param6,
                             model.param7, model.param8, model.param9).detach().numpy().flatten()

# 确保所有数据为numpy类型
    true_values = targets_data.flatten()[:]
    predicted_values = np.array(predicted_values)
    # 获取对比图图像
    comparison_buf = plot_comparison(partial_time, true_values, predicted_values, physic_values)


# 损失曲线
def plot_loss_curves(losses, data_metrics, physics_metrics):
    plt.figure()
    plt.plot(np.log10(losses), label='Total Loss')
    plt.plot(np.log10(data_metrics), label='Data Loss')
    plt.plot(np.log10(physics_metrics), label='Physics Loss')
    plt.xlabel('Epoch')
    plt.ylabel('Log Loss')
    plt.legend()
    plt.grid()
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight')
    plt.close()
    buf.seek(0)
    return buf

# 获取损失曲线图图像
loss_buf = plot_loss_curves(losses, data_metrics, physics_metrics)

# 参数输出
# 捕获print输出
original_stdout = sys.stdout
sys.stdout = string_io = io.StringIO()
print(f"关键参数:\n"
      f"T_initial={model.T_initial.item():.1f} K\n"
      f"param1={model.param1.item():.3f}\n"
      f"param2={model.param2.item():.3f}\n"
      f"param3={model.param3.item():.3f}\n"
      f"param4={model.param4.item():.5f}\n"
      f"param5={model.param5.item():.3f}\n"
      f"param6={model.param6.item():.3f}\n"
      f"param7={model.param7.item():.3f}")
sys.stdout = original_stdout
param_text = string_io.getvalue()

# 创建PPT并插入内容
def create_ppt(temp_buf, comparison_buf, loss_buf, param_text, output_path=r'E:\Desktop\results.pptx'):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_slide_layout)

    # 插入温度图（左上）
    left = Inches(0.5)
    top = Inches(0.5)
    slide.shapes.add_picture(temp_buf, left, top, width=Inches(4), height=Inches(2.5))

    # 插入对比图（右上）
    left = Inches(5)
    slide.shapes.add_picture(comparison_buf, left, top, width=Inches(4), height=Inches(2.5))

    # 插入损失曲线图（左下）
    top = Inches(3.5)
    left = Inches(0.5)
    slide.shapes.add_picture(loss_buf, left, top, width=Inches(4), height=Inches(2.5))

    # 插入参数文本框（右下）
    left = Inches(5)
    width = Inches(4)
    height = Inches(2.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = param_text

    # 设置文本框样式
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.size = Pt(12)  # 12pt
        paragraph.font.name = '宋体'
        paragraph.font.color.rgb = RGBColor(0, 0, 0)

    prs.save(output_path)
    # 新增功能：自动打开生成的PPT
    if os.name == 'nt':  # Windows系统
        os.startfile(output_path)

# 调用函数生成PPT
create_ppt(temp_buf, comparison_buf, loss_buf, param_text)

# # 7. 温度变化可视化
# def plot_temperature(time, temperature):
# plt.figure(figsize=(10, 6))
# plt.plot(time, temperature, label="Temperature", color='r', linewidth=2)
# plt.xlabel("Time (s)")
# plt.ylabel("Temperature")
# plt.title("Device Temperature Evolution")
# plt.legend()
# plt.grid(True)
# plt.show()

# # 获取完整温度序列（使用全量数据）
# with torch.no_grad():
# # 准备完整数据
# full_inputs = torch.tensor(X_scaled, dtype=torch.float32)
# full_time = full_inputs[:, 0] * input_std[0] + input_mean[0]
# full_vds = full_inputs[:, 1] * input_std[1] + input_mean[1]
# full_vgs = full_inputs[:, 2] * input_std[2] + input_mean[2]

# # 计算温度序列
# T_sequence, _ = integrate_temperature(full_time, full_vds, full_vgs, model, physics_model)

# # 绘制温度变化
# plot_temperature(full_time.numpy(), T_sequence.numpy())

# # 对比结果可视化（使用前300个样本）
# def plot_comparison(time, true, pred, physics):
# plt.figure(figsize=(10, 6))
# plt.plot(time, true, label="True Values", marker='o', markersize=3)
# plt.plot(time, pred, label="Predicted Values", marker='x', markersize=3)
# plt.plot(time, physics, label="Physics Values", marker='.', markersize=3)
# plt.xlabel("Time (s)")
# plt.ylabel("Drain Current (A)")
# plt.title("Current Comparison Over Time")
# plt.legend()
# plt.grid(True)
# plt.show()

# # 主可视化流程
# # 修复AttributeError的核心修改部分
# with torch.no_grad():
# # 2. 准备对比数据（前300个样本）
# # 正确获取时间序列（全numpy运算）
# input_std_np = input_std.numpy()
# input_mean_np = input_mean.numpy()
# partial_time = X_scaled[:, 0] * input_std_np[0] + input_mean_np[0]

# # 修复神经网络预测值的转换
# predicted_values = [neural(Vds, Vgs, T) for Vds, Vgs, T in inputs_data[:]] # 原始输出类型

# # 正确转换物理模型输出
# # 修复部分：批量计算物理模型的预测值 确保physic_values为numpy类型
# physic_values = physics_model(full_vds, full_vgs, T_sequence,
# model.param1, model.param2, model.param3,
# model.param4, model.param5, model.param6,
# model.param7, model.param8, model.param9).detach().numpy().flatten()

# # 确保所有数据为numpy类型
# true_values = targets_data.flatten()[:]
# predicted_values = np.array(predicted_values)
# plot_comparison(partial_time, true_values, predicted_values, physic_values)

# # 其他可视化保持不变
# # 损失曲线
# plt.figure()
# plt.plot(np.log10(losses), label='Total Loss')
# plt.plot(np.log10(data_metrics), label='Data Loss')
# plt.plot(np.log10(physics_metrics), label='Physics Loss')
# plt.xlabel('Epoch'), plt.ylabel('Log Loss'), plt.legend(), plt.grid()
# plt.show()

# # 参数输出
# print(f"关键参数:\n"
# f"T_initial={model.T_initial.item():.1f} K\n"
# f"param1={model.param1.item():.3f}\n"
# f"param2={model.param2.item():.3f}\n"
# f"param3={model.param3.item():.3f}\n"
# f"param4={model.param4.item():.5f}\n"
# f"param5={model.param5.item():.3f}\n"
# f"param6={model.param6.item():.3f}\n"
# f"param7={model.param7.item():.3f}"
# )
